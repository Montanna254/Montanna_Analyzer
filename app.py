import streamlit as st

def login():
    password = st.text_input("Enter password", type="password")
    if st.button("Login"):
        if password == "YourPassword":
            st.session_state["password_correct"] = True
        else:
            st.session_state["password_correct"] = False
            st.error("Incorrect password")

if "password_correct" not in st.session_state:
    login()

if st.session_state.get("password_correct", False):
    st.success("Welcome!")
    # The rest of your Streamlit app goes here

import os
import json
from typing import List, Dict
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
import numpy as np
import streamlit as st
from io import StringIO

# Text extraction tools
import docx
import fitz  # PyMuPDF
from pptx import Presentation
from PIL import Image
import pytesseract
import requests
from bs4 import BeautifulSoup

class ArticleAnalyzer:
    def __init__(self):
        self.articles: List[str] = []
        self.vectorizer = TfidfVectorizer(stop_words='english')
        self.article_vectors = None

    def load_articles(self, article_list: List[str]):
        if not article_list:
            raise ValueError("No articles provided for analysis.")
        self.articles = article_list
        self.article_vectors = self.vectorizer.fit_transform(self.articles)

    def analyze_prompt(self, prompt: str, top_n: int = 5) -> List[Dict[str, any]]:
        if self.article_vectors is None:
            raise RuntimeError("No articles loaded. Please load articles before analysis.")

        prompt_vector = self.vectorizer.transform([prompt])
        similarities = cosine_similarity(prompt_vector, self.article_vectors).flatten()

        top_indices = np.argsort(similarities)[-top_n:][::-1]
        results = [
            {
                "index": i,
                "similarity_score": float(similarities[i]),
                "article_excerpt": self.articles[i][:300]
            }
            for i in top_indices
        ]
        return results

# Extraction functions
def extract_text_from_docx(file) -> str:
    doc = docx.Document(file)
    return "\n".join([p.text for p in doc.paragraphs])

def extract_text_from_pdf(file) -> str:
    doc = fitz.open(stream=file.read(), filetype="pdf")
    return "\n".join([page.get_text() for page in doc])

def extract_text_from_pptx(file) -> str:
    prs = Presentation(file)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def extract_text_from_image(file) -> str:
    image = Image.open(file)
    return pytesseract.image_to_string(image)

def fetch_articles_from_web(keyword: str, num_articles: int = 5) -> List[str]:
    headers = {'User-Agent': 'Mozilla/5.0'}
    query = keyword.replace(' ', '+')
    search_url = f"https://www.google.com/search?q={query}"

    response = requests.get(search_url, headers=headers)
    soup = BeautifulSoup(response.text, "html.parser")
    links = [a['href'] for a in soup.select('a') if 'url?q=' in a['href']]

    articles = []
    for link in links[:num_articles]:
        try:
            real_url = link.split('url?q=')[1].split('&')[0]
            article_html = requests.get(real_url, headers=headers)
            article_soup = BeautifulSoup(article_html.text, "html.parser")
            text = ' '.join([p.get_text() for p in article_soup.find_all('p')])
            if text.strip():
                articles.append(text)
        except:
            continue
    return articles

# Streamlit app
st.title("Montana Article Analyzer")
st.write("Upload one or more documents (.txt, .pdf, .docx, .pptx, .png, .jpg) or search for web articles.")

uploaded_files = st.file_uploader("Choose your documents", type=["txt", "pdf", "docx", "pptx", "png", "jpg"], accept_multiple_files=True)
prompt = st.text_input("Enter your analysis prompt:")
top_n = st.slider("Number of top matching articles to return", 1, 50, 5)

st.subheader("Search the Web")
keyword = st.text_input("Enter keyword to search for articles on the web:")
num_online_articles = st.slider("Number of web articles to fetch", 1, 10, 3)

article_texts = []

if uploaded_files:
    for file in uploaded_files:
        filename = file.name.lower()
        try:
            if filename.endswith(".txt"):
                text = StringIO(file.read().decode("utf-8")).read()
            elif filename.endswith(".pdf"):
                text = extract_text_from_pdf(file)
            elif filename.endswith(".docx"):
                text = extract_text_from_docx(file)
            elif filename.endswith(".pptx"):
                text = extract_text_from_pptx(file)
            elif filename.endswith((".png", ".jpg")):
                text = extract_text_from_image(file)
            else:
                text = ""
            if text.strip():
                article_texts.append(text)
        except Exception as e:
            st.error(f"Error reading {filename}: {e}")

if keyword:
    try:
        web_articles = fetch_articles_from_web(keyword, num_online_articles)
        article_texts.extend(web_articles)
        st.success(f"Fetched {len(web_articles)} articles from the web.")
    except Exception as e:
        st.error(f"Error fetching web articles: {e}")

if article_texts and prompt:
    analyzer = ArticleAnalyzer()
    analyzer.load_articles(article_texts)
    results = analyzer.analyze_prompt(prompt, top_n=top_n)

    st.subheader("Top Matching Articles:")
    for result in results:
        st.markdown(f"**Article {result['index']} (Score: {result['similarity_score']:.2f})**")
        st.write(result['article_excerpt'])
elif prompt:
    st.warning("Please upload documents or enter a keyword to search for web articles.")
